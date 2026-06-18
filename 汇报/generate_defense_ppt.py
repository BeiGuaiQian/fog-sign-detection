from __future__ import annotations

from pathlib import Path

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


REPORT = Path(__file__).resolve().parent
ROOT = REPORT.parent
TEMPLATE = REPORT / "32.pptx"
OUT = REPORT / "雾天路标检测系统_8分钟答辩稿.pptx"
LOGO_SRC = REPORT / "logo800_16491858436941550.jpg"
EMBLEM = REPORT / "ecnu_emblem_only.png"

FOGGY_SHEET = ROOT / "data" / "foggy_test" / "foggy_contact_sheet.jpg"
TRAIN_DIR = ROOT.parent / "training-runs" / "traffic_sign_yolov8s_960-2"
RESULTS_CSV = TRAIN_DIR / "results.csv"
VAL_PRED = TRAIN_DIR / "val_batch0_pred.jpg"

RED = RGBColor(171, 31, 53)
DARK = RGBColor(31, 42, 68)
MID = RGBColor(92, 99, 112)
LIGHT_BG = RGBColor(248, 249, 252)
PANEL = RGBColor(255, 255, 255)
BORDER = RGBColor(224, 228, 236)
BLUE = RGBColor(49, 83, 129)
GREEN = RGBColor(69, 129, 93)
GOLD = RGBColor(196, 151, 72)
PALE_RED = RGBColor(253, 247, 249)
PALE_RED_LINE = RGBColor(238, 209, 216)

FONT = "Microsoft YaHei"
FONT_BOLD = "Microsoft YaHei UI"


def crop_emblem() -> None:
    """Create a transparent image containing only the ECNU round emblem."""
    im = Image.open(LOGO_SRC).convert("RGBA")
    region = im.crop((0, 0, 360, im.height))
    xs, ys = [], []
    for y in range(region.height):
        for x in range(region.width):
            r, g, b, a = region.getpixel((x, y))
            if a and not (r > 245 and g > 245 and b > 245):
                xs.append(x)
                ys.append(y)

    if xs:
        box = (
            max(0, min(xs) - 16),
            max(0, min(ys) - 16),
            min(region.width, max(xs) + 16),
            min(region.height, max(ys) + 16),
        )
    else:
        box = (55, 300, 315, 575)

    crop = region.crop(box)
    side = max(crop.width, crop.height)
    canvas = Image.new("RGBA", (side, side), (255, 255, 255, 0))
    canvas.paste(crop, ((side - crop.width) // 2, (side - crop.height) // 2))

    for y in range(side):
        for x in range(side):
            r, g, b, _ = canvas.getpixel((x, y))
            if r > 248 and g > 248 and b > 248:
                canvas.putpixel((x, y), (255, 255, 255, 0))

    canvas.save(EMBLEM)


def solid(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def textbox(slide, content, x, y, w, h, size=14, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    para = frame.paragraphs[0]
    para.text = content
    para.alignment = align
    para.font.name = FONT_BOLD if bold else FONT
    para.font.size = Pt(size)
    para.font.bold = bold
    para.font.color.rgb = color
    return box


def add_background(slide, prs) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    solid(bg, LIGHT_BG)
    bg.line.fill.background()

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.58))
    solid(band, RGBColor(255, 255, 255))
    band.line.color.rgb = BORDER

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.58), prs.slide_width, Inches(0.035))
    solid(line, RED)
    line.line.fill.background()


def add_header(slide, prs, section, index) -> None:
    add_background(slide, prs)
    slide.shapes.add_picture(str(EMBLEM), Inches(0.28), Inches(0.12), width=Inches(0.36), height=Inches(0.36))
    textbox(slide, section, 0.78, 0.11, 10.7, 0.36, size=17, color=DARK, bold=True)
    textbox(slide, f"{index:02d}", 12.1, 0.16, 0.8, 0.3, size=10, color=MID, align=PP_ALIGN.RIGHT)


def add_title(slide, title, x=0.75, y=0.92, w=10, h=0.55, size=25):
    return textbox(slide, title, x, y, w, h, size=size, color=DARK, bold=True)


def add_card(slide, x, y, w, h, fill=PANEL, line=BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_bullets(slide, items, x, y, w, h, size=14, color=DARK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for i, item in enumerate(items):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.text = item
        para.level = 0
        para.font.name = FONT
        para.font.size = Pt(size)
        para.font.color.rgb = color
        para.space_after = Pt(7)
        para._p.get_or_add_pPr().set("marL", "228600")
        para._p.get_or_add_pPr().set("hanging", "91440")
    return box


def add_label(slide, content, x, y, w=0.58, color=RED):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.32))
    solid(shape, color)
    shape.line.fill.background()
    textbox(slide, content, x, y + 0.05, w, 0.22, size=9, color=RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)


def add_image_fit(slide, image_path: Path, x, y, w, h, border=True):
    image = Image.open(image_path)
    iw, ih = image.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        width = w
        height = w / img_ratio
        left = x
        top = y + (h - height) / 2
    else:
        height = h
        width = h * img_ratio
        left = x + (w - width) / 2
        top = y
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))
    if border:
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        frame.fill.background()
        frame.line.color.rgb = BORDER
        frame.line.width = Pt(1)


def add_flow(slide, steps, x, y, w, h):
    gap = 0.18
    box_width = (w - gap * (len(steps) - 1)) / len(steps)
    colors = [RED, BLUE, GREEN, GOLD, DARK]
    for i, step in enumerate(steps):
        bx = x + i * (box_width + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bx), Inches(y), Inches(box_width), Inches(h))
        solid(shape, colors[i % len(colors)])
        shape.line.fill.background()
        textbox(slide, step, bx + 0.05, y + 0.13, box_width - 0.1, h - 0.2, size=13, color=RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(bx + box_width + 0.035), Inches(y + 0.25), Inches(0.12), Inches(0.22))
            solid(arrow, MID)
            arrow.line.fill.background()


def load_metrics():
    default = {"precision": "0.960", "recall": "0.946", "map50": "0.963", "map5095": "0.785", "epoch": "118"}
    if not RESULTS_CSV.exists():
        return default
    data = pd.read_csv(RESULTS_CSV)
    row = data.iloc[-1]
    return {
        "precision": f"{row['metrics/precision(B)']:.3f}",
        "recall": f"{row['metrics/recall(B)']:.3f}",
        "map50": f"{row['metrics/mAP50(B)']:.3f}",
        "map5095": f"{row['metrics/mAP50-95(B)']:.3f}",
        "epoch": str(int(row["epoch"])),
    }


def normalize_fonts(prs) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = FONT


def build() -> None:
    crop_emblem()
    metrics = load_metrics()

    prs = Presentation(TEMPLATE)
    for slide_id in list(prs.slides._sldIdLst):
        prs.part.drop_rel(slide_id.rId)
        prs.slides._sldIdLst.remove(slide_id)

    blank = prs.slide_layouts[1]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    add_background(slide, prs)
    slide.shapes.add_picture(str(EMBLEM), Inches(0.42), Inches(0.28), width=Inches(0.46), height=Inches(0.46))
    textbox(slide, "雾天交通视觉 · 图像增强 · 目标定位", 0.98, 0.34, 5.8, 0.3, size=12, color=MID)
    textbox(slide, "基于图像去雾与目标定位的\n雾天路标检测系统", 0.82, 1.55, 7.2, 1.35, size=34, color=DARK, bold=True)
    textbox(slide, "面向雾天场景，构建“去雾处理—路标检测—指标评价—可视化展示”的完整实验流程。", 0.86, 3.12, 7.1, 0.55, size=15, color=MID)
    add_flow(slide, ["雾天图像", "DCP 去雾", "YOLOv8 检测", "指标对比", "系统展示"], 0.9, 4.35, 7.5, 0.72)
    if FOGGY_SHEET.exists():
        add_image_fit(slide, FOGGY_SHEET, 8.75, 1.15, 3.8, 5.45)
    textbox(slide, "课程项目答辩", 0.9, 6.78, 2, 0.3, size=11, color=MID)

    # 2. Background
    slide = prs.slides.add_slide(blank)
    add_header(slide, prs, "背景与问题", 2)
    add_title(slide, "雾天会削弱路标的可见线索")
    textbox(slide, "路标识别依赖颜色、形状和边缘。雾天会同时压低这些信息，检测模型的输入质量随之下降。", 0.78, 1.48, 10.6, 0.42, size=13, color=MID)
    for i, (title, desc, color) in enumerate([
        ("对比度下降", "图像整体发白，目标边界变弱。", RED),
        ("颜色辨识度降低", "红灯、绿灯、警告牌等类别更易混淆。", BLUE),
        ("小目标细节丢失", "远距离路标在雾中只保留少量有效像素。", GREEN),
    ]):
        x = 0.82 + i * 3.95
        add_card(slide, x, 2.16, 3.55, 2.3)
        add_label(slide, f"0{i + 1}", x + 0.22, 2.38, color=color)
        textbox(slide, title, x + 0.28, 2.86, 2.95, 0.34, size=18, bold=True)
        textbox(slide, desc, x + 0.28, 3.34, 2.88, 0.72, size=13, color=MID)
    add_card(slide, 1.15, 5.15, 11.05, 0.95, fill=PALE_RED, line=PALE_RED_LINE)
    textbox(slide, "本项目关注的问题：通过去雾增强改善图像可见性，并比较处理前后路标检测结果的变化。", 1.45, 5.44, 10.4, 0.36, size=16, color=RED, bold=True)

    # 3. Workflow
    slide = prs.slides.add_slide(blank)
    add_header(slide, prs, "系统总体流程", 3)
    add_title(slide, "去雾、检测、评价三步串起来")
    add_flow(slide, ["输入雾天图像", "DCP 图像去雾", "YOLOv8 双路检测", "指标对比", "页面展示"], 0.78, 1.95, 11.65, 0.86)
    for i, (title, desc) in enumerate([
        ("图像去雾", "输出去雾图、暗通道图和透射率图，保留算法中间结果。"),
        ("目标检测", "原始雾图和去雾图使用同一检测模型，结果具有可比性。"),
        ("指标评价", "同时统计图像质量和检测结果，形成可解释的对比。"),
        ("交互展示", "通过页面完成上传、调参、观察、下载，便于课堂演示。"),
    ]):
        row, col = divmod(i, 2)
        x, y = 0.85 + col * 5.9, 3.28 + row * 1.25
        add_card(slide, x, y, 5.35, 0.95)
        textbox(slide, title, x + 0.22, y + 0.17, 1.4, 0.28, size=14, color=RED, bold=True)
        textbox(slide, desc, x + 1.52, y + 0.16, 3.55, 0.42, size=12.5)

    # 4. DCP
    slide = prs.slides.add_slide(blank)
    add_header(slide, prs, "DCP 去雾方法", 4)
    add_title(slide, "用暗通道估计雾的影响")
    add_card(slide, 0.82, 1.62, 5.15, 4.95)
    textbox(slide, "大气散射模型", 1.08, 1.9, 3, 0.32, size=16, color=RED, bold=True)
    textbox(slide, "I(x) = J(x)t(x) + A(1 - t(x))", 1.08, 2.42, 4.4, 0.4, size=20, bold=True)
    add_bullets(slide, ["I(x)：输入的有雾图像", "J(x)：希望恢复的清晰图像", "A：雾天环境中的大气光", "t(x)：场景信息保留下来的比例"], 1.06, 3.08, 4.5, 1.65, size=13)
    textbox(slide, "实现流程：暗通道计算 → 大气光估计 → 透射率估计 → 引导滤波 → 图像恢复", 1.08, 5.33, 4.55, 0.6, size=13, color=MID)
    for i, item in enumerate(["暗通道", "透射率", "去雾图"]):
        x = 6.55 + i * 1.85
        add_card(slide, x, 2.1, 1.45, 1.1)
        textbox(slide, item, x + 0.1, 2.47, 1.25, 0.28, size=14, color=[RED, BLUE, GREEN][i], bold=True, align=PP_ALIGN.CENTER)
        if i < 2:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 1.52), Inches(2.47), Inches(0.22), Inches(0.18))
            solid(arrow, MID)
            arrow.line.fill.background()
    add_bullets(slide, ["暗通道反映雾对局部区域的抬亮效果。", "透射率图描述不同位置的雾浓度差异。", "引导滤波让恢复结果更平滑、边缘更自然。"], 6.42, 3.72, 5.55, 1.48, size=14)
    textbox(slide, "这一模块提供去雾结果，也提供暗通道和透射率图作为可解释的中间结果。", 6.45, 5.78, 5.5, 0.45, size=13, color=MID)

    # 5. YOLO
    slide = prs.slides.add_slide(blank)
    add_header(slide, prs, "YOLOv8 路标检测", 5)
    add_title(slide, "对原图和去雾图做同一套检测")
    add_bullets(slide, ["检测框架采用 Ultralytics YOLOv8。", "原始雾图和去雾图使用相同权重、相同阈值进行推理。", "输出类别、置信度、检测框坐标和目标面积。", "检测结果以图像和表格两种方式呈现。"], 0.9, 1.75, 5.25, 2.0, size=15)
    add_card(slide, 0.9, 4.42, 5.15, 1.28, fill=PALE_RED, line=PALE_RED_LINE)
    textbox(slide, "对比设计", 1.18, 4.68, 1.4, 0.28, size=15, color=RED, bold=True)
    textbox(slide, "同一张图像分别经过“直接检测”和“去雾后检测”，结果差异来自图像增强过程。", 2.18, 4.66, 3.55, 0.44, size=12.5)
    if VAL_PRED.exists():
        add_image_fit(slide, VAL_PRED, 6.55, 1.45, 5.6, 4.8)
        textbox(slide, "检测结果示例：模型输出路标位置、类别和置信度。", 6.75, 6.32, 5.2, 0.25, size=11, color=MID)

    # 6. Dataset
    slide = prs.slides.add_slide(blank)
    add_header(slide, prs, "数据集处理与训练设置", 6)
    add_title(slide, "让每个类别都进入训练、验证和测试")
    textbox(slide, "目标检测数据可能一张图含多个类别，因此采用多标签分层划分来保持类别覆盖。", 0.78, 1.48, 10.6, 0.35, size=13, color=MID)
    for i, (name, value, color) in enumerate([("训练集", "1674 张", RED), ("验证集", "209 张", BLUE), ("测试集", "210 张", GREEN)]):
        x = 0.95 + i * 3.9
        add_card(slide, x, 2.05, 3.2, 1.35)
        textbox(slide, name, x + 0.25, 2.28, 1.1, 0.3, size=14, color=MID)
        textbox(slide, value, x + 0.25, 2.72, 2, 0.38, size=24, color=color, bold=True)
    add_bullets(slide, ["数据集包含 21 类交通标志。", "每个类别均覆盖训练、验证和测试集合。", "训练使用 YOLOv8s，输入尺寸为 960，适配小目标场景。", "验证集指标用于观察模型整体收敛和类别混淆情况。"], 1.0, 4.05, 5.2, 1.78, size=14)
    for i, (label, value) in enumerate([("Precision", metrics["precision"]), ("Recall", metrics["recall"]), ("mAP50", metrics["map50"]), ("mAP50-95", metrics["map5095"])]):
        x, y = 6.72 + (i % 2) * 2.45, 4.0 + (i // 2) * 1.02
        add_card(slide, x, y, 2.08, 0.78)
        textbox(slide, label, x + 0.16, y + 0.12, 0.85, 0.22, size=10.5, color=MID)
        textbox(slide, value, x + 1.1, y + 0.1, 0.72, 0.28, size=18, color=RED if i in (0, 2) else BLUE, bold=True, align=PP_ALIGN.RIGHT)
    textbox(slide, f"训练记录最后一轮：epoch {metrics['epoch']}", 6.78, 6.07, 3.8, 0.3, size=11, color=MID)

    # 7. Streamlit system
    slide = prs.slides.add_slide(blank)
    add_header(slide, prs, "可视化系统", 7)
    add_title(slide, "把算法过程做成可操作的实验页面")
    add_bullets(slide, ["上传一张雾天路标图像即可运行完整流程。", "侧边栏支持调节去雾参数、置信度阈值和推理设备。", "页面同步展示原始雾图、去雾图、暗通道图、透射率图。", "检测框表格和指标表用于解释结果变化。", "去雾图和检测结果图可以直接下载。"], 0.9, 1.75, 5.15, 2.55, size=14.5)
    add_card(slide, 6.45, 1.45, 5.75, 4.95)
    add_card(slide, 6.75, 1.78, 1.25, 4.1, fill=RGBColor(242, 244, 248))
    textbox(slide, "参数区", 6.91, 2.02, 0.9, 0.3, size=12, color=RED, bold=True, align=PP_ALIGN.CENTER)
    for j, item in enumerate(["DCP", "YOLO", "设备", "阈值"]):
        add_card(slide, 6.9, 2.55 + j * 0.65, 0.9, 0.35)
        textbox(slide, item, 6.96, 2.62 + j * 0.65, 0.78, 0.16, size=8.5, color=MID, align=PP_ALIGN.CENTER)
    add_card(slide, 8.18, 1.78, 3.72, 1.6)
    textbox(slide, "图像展示区", 9.3, 2.38, 1.4, 0.28, size=13, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_card(slide, 8.18, 3.62, 3.72, 0.92)
    textbox(slide, "检测框信息表", 9.15, 3.95, 1.75, 0.25, size=12, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_card(slide, 8.18, 4.75, 3.72, 1.13)
    textbox(slide, "指标对比与下载", 9.08, 5.2, 1.9, 0.25, size=12, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

    # 8. Metrics
    slide = prs.slides.add_slide(blank)
    add_header(slide, prs, "评价指标设计", 8)
    add_title(slide, "同时看图像质量和检测表现")
    textbox(slide, "同一张图像在去雾前后进行对比，图像指标和检测指标共同说明效果变化。", 0.78, 1.48, 10.6, 0.35, size=13, color=MID)
    add_card(slide, 0.92, 2.05, 5.45, 3.85)
    textbox(slide, "图像质量指标", 1.2, 2.32, 2, 0.32, size=17, color=RED, bold=True)
    add_bullets(slide, ["对比度：观察灰度标准差变化。", "信息熵：描述图像信息量。", "清晰度：使用拉普拉斯方差衡量边缘强度。", "边缘数量：反映可见结构是否增加。"], 1.18, 2.9, 4.6, 1.75, size=13.5)
    add_card(slide, 6.85, 2.05, 5.45, 3.85)
    textbox(slide, "检测结果指标", 7.13, 2.32, 2, 0.32, size=17, color=BLUE, bold=True)
    add_bullets(slide, ["检测数量：统计模型响应的目标个数。", "平均置信度：观察整体预测把握程度。", "最高置信度：观察最稳定目标的预测强度。", "检测框面积：辅助分析小目标和大目标差异。"], 7.1, 2.9, 4.6, 1.75, size=13.5)
    textbox(slide, "指标解释原则：图像更清晰、模型响应更稳定、检测结果更容易分析。", 2.15, 6.28, 9.1, 0.32, size=14, bold=True, align=PP_ALIGN.CENTER)

    # 9. Results
    slide = prs.slides.add_slide(blank)
    add_header(slide, prs, "实验展示与结果分析", 9)
    add_title(slide, "去雾后可见性提升，检测结果随场景变化")
    add_bullets(slide, ["测试样例覆盖信号灯、警告牌、禁令牌、人行横道、公交站等类别。", "DCP 去雾通常提升局部对比度和边缘可见性。", "模型置信度和检测框稳定性会随目标尺寸、雾浓度和颜色保持程度变化。"], 0.9, 1.55, 4.85, 1.55, size=13.5)
    if FOGGY_SHEET.exists():
        add_image_fit(slide, FOGGY_SHEET, 6.05, 1.25, 6.15, 5.65)
    add_card(slide, 0.9, 3.55, 4.85, 2.45, fill=PALE_RED, line=PALE_RED_LINE)
    textbox(slide, "观察结论", 1.18, 3.82, 1.3, 0.3, size=15, color=RED, bold=True)
    add_bullets(slide, ["去雾改善了路标轮廓和局部纹理。", "颜色敏感目标仍需要关注色彩偏移。", "小目标在重雾中仍然是主要难点。"], 1.18, 4.28, 4.2, 1.0, size=12.5)

    # 10. Limitations and future work
    slide = prs.slides.add_slide(blank)
    add_header(slide, prs, "不足与改进方向", 10)
    add_title(slide, "系统已经可用，后续还能做得更稳")
    textbox(slide, "当前边界", 1.0, 1.75, 1.4, 0.3, size=16, color=RED, bold=True)
    textbox(slide, "改进方向", 7.05, 1.75, 1.4, 0.3, size=16, color=BLUE, bold=True)
    for i, (title, desc) in enumerate([
        ("DCP 稳定性", "天空、强光和大面积白色区域会影响大气光估计。"),
        ("颜色保持", "去雾可能带来颜色偏移，信号灯等类别更敏感。"),
        ("小目标检测", "远距离路标在重雾中仍然保留较少有效细节。"),
    ]):
        y = 2.22 + i * 1.18
        add_card(slide, 0.92, y, 5.25, 0.82)
        textbox(slide, title, 1.18, y + 0.14, 1.5, 0.24, size=13, color=RED, bold=True)
        textbox(slide, desc, 2.52, y + 0.13, 3.2, 0.3, size=11.5)
    for i, (title, desc) in enumerate([
        ("学习型去雾", "加入 AOD-Net、DehazeFormer 等方法进行对比。"),
        ("批量评测", "增加页面端批量上传、指标汇总和报告导出。"),
        ("标准指标", "继续引入 Precision、Recall、mAP 和混淆矩阵分析。"),
    ]):
        y = 2.22 + i * 1.18
        add_card(slide, 6.82, y, 5.25, 0.82)
        textbox(slide, title, 7.08, y + 0.14, 1.5, 0.24, size=13, color=BLUE, bold=True)
        textbox(slide, desc, 8.42, y + 0.13, 3.25, 0.3, size=11.5)
    add_card(slide, 2.12, 5.88, 9.1, 0.66, fill=RGBColor(250, 250, 252))
    textbox(slide, "下一步重点：把单图演示扩展为批量实验平台，并用更完整的检测指标支撑结论。", 2.45, 6.07, 8.45, 0.22, size=13, bold=True, align=PP_ALIGN.CENTER)

    # 11. Summary
    slide = prs.slides.add_slide(blank)
    add_header(slide, prs, "总结", 11)
    add_title(slide, "完成了一个雾天路标检测实验系统")
    for i, (title, desc, color) in enumerate([
        ("图像去雾", "实现 DCP 暗通道先验，并输出可解释中间结果。", RED),
        ("目标检测", "接入 YOLOv8 路标检测，支持自定义训练权重。", BLUE),
        ("指标评价", "同时比较图像质量和检测结果，支撑实验分析。", GREEN),
        ("交互展示", "通过可视化页面完成上传、调参、展示和下载。", GOLD),
    ]):
        x = 0.95 + (i % 2) * 5.9
        y = 1.92 + (i // 2) * 1.35
        add_card(slide, x, y, 5.2, 0.95)
        add_label(slide, f"{i + 1}", x + 0.22, y + 0.2, w=0.48, color=color)
        textbox(slide, title, x + 0.9, y + 0.18, 1.3, 0.25, size=14, color=color, bold=True)
        textbox(slide, desc, x + 2.08, y + 0.17, 2.85, 0.32, size=11.5)
    add_card(slide, 1.18, 5.15, 10.95, 0.95, fill=PALE_RED, line=PALE_RED_LINE)
    textbox(slide, "本系统把图像增强、目标检测、指标评价和交互展示整合到同一条实验流程中，用于观察雾天图像处理对路标检测的影响。", 1.48, 5.43, 10.3, 0.36, size=15, color=RED, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, "以上是我的项目汇报，请老师批评指正。", 4.35, 6.58, 4.7, 0.3, size=14, color=MID, align=PP_ALIGN.CENTER)

    normalize_fonts(prs)
    prs.save(OUT)
    print(OUT)
    print(f"slides={len(prs.slides)}")


if __name__ == "__main__":
    build()
